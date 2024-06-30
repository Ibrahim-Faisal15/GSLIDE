from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from text_gen import get_response
from image_gen import get_image
from io import BytesIO




def get_presentation(topic):
    slide1_title, slide2_content, slide_t, slide_c = get_response(topic)
    presentation = Presentation()
    ppt_bytes = BytesIO()
    get_image(topic)
    print("here is the slide " ,  slide_t)
    print(len(slide_t))

    slide0 = presentation.slides.add_slide(presentation.slide_layouts[0])
    title0 = slide0.shapes.title
    title0.text = slide1_title
    title0.text_frame.paragraphs[0].runs[0].font.bold = True

    subtitle_placeholder = slide0.placeholders[1]
    subtitle_placeholder.text = ""

    slide0.background.fill.solid()
    slide0.background.fill.fore_color.rgb = RGBColor(0, 0, 0)

    for shape in slide0.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)

    title0.text_frame.paragraphs[0].runs[0].font.size = Pt(80)
    # Slide 1
    slide1 = presentation.slides.add_slide(presentation.slide_layouts[3])
    title1 = slide1.shapes.title
    title1.text = slide_t[0]
    title1.text_frame.paragraphs[0].runs[0].font.bold = True

    content1 = slide1.placeholders[1]
    content1.text = slide_c[0]

    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
    #
    for shape in slide1.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.font.size = Pt(18)
    title1.text_frame.paragraphs[0].runs[0].font.size = Pt(40)

    image_placeholder = slide1.placeholders[2]
    for shape in slide1.shapes:
        if shape == image_placeholder:
            slide1.shapes._spTree.remove(shape._element)

    # Add the image to the placeholder
    image_path1 = 'images/1.jpg'
    left1 = image_placeholder.left
    top1 = image_placeholder.top
    width1 = image_placeholder.width
    height1 = image_placeholder.height

    # Calculate the aspect ratio of the image
    image_ratio = width1 / height1

    # Load the image to get its original dimensions
    image = slide1.shapes.add_picture(image_path1, left1, top1)
    original_width = image.width
    original_height = image.height

    # Calculate the new dimensions to fit within the placeholder
    if image_ratio > 1:
        new_width = width1
        new_height = int(new_width / image_ratio)
    else:
        new_height = height1
        new_width = int(new_height * image_ratio)

    # Calculate the offset to center the image within the placeholder
    left_offset = (width1 - new_width) // 2
    top_offset = (height1 - new_height) // 2

    # Remove the temporary image shape
    slide1.shapes._spTree.remove(image._element)

    # Add the resized image to the placeholder
    slide1.shapes.add_picture(image_path1, left1 + left_offset, top1 + top_offset, new_width, new_height)

    # Slide 2
    slide2 = presentation.slides.add_slide(presentation.slide_layouts[3])
    title2 = slide2.shapes.title
    title2.text = slide_t[1]
    content2 = slide2.placeholders[1]
    content2.text = slide_c[1]

    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
    #
    for shape in slide2.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # RGB values for white
                    run.font.size = Pt(18)

    title2.text_frame.paragraphs[0].runs[0].font.size = Pt(40)
    title2.text_frame.paragraphs[0].runs[0].font.bold = True

    image_placeholder = slide2.placeholders[2]
    for shape in slide2.shapes:
        if shape == image_placeholder:
            slide2.shapes._spTree.remove(shape._element)

    image_path2 = 'images/2.jpg'
    left2 = image_placeholder.left
    top2 = image_placeholder.top
    width2 = image_placeholder.width
    height2 = image_placeholder.height

    # Calculate the aspect ratio of the image
    image_ratio = width2 / height2

    # Load the image to get its original dimensions
    image = slide2.shapes.add_picture(image_path2, left2, top2)
    original_width = image.width
    original_height = image.height

    # Calculate the new dimensions to fit within the placeholder
    if image_ratio > 1:
        new_width = width2
        new_height = int(new_width / image_ratio)
    else:
        new_height = height2
        new_width = int(new_height * image_ratio)

    # Calculate the offset to center the image within the placeholder
    left_offset = (width2 - new_width) // 2
    top_offset = (height2 - new_height) // 2

    # Remove the temporary image shape
    slide2.shapes._spTree.remove(image._element)

    # Add the resized image to the placeholder
    slide2.shapes.add_picture(image_path2, left1 + left_offset, top1 + top_offset, new_width, new_height)

    slide3 = presentation.slides.add_slide(presentation.slide_layouts[3])
    title3 = slide3.shapes.title
    title3.text = slide_t[0]
    content3 = slide3.placeholders[1]
    content3.text = slide_c[0]

    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
    #
    for shape in slide3.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # RGB values for white
                    run.font.size = Pt(18)

    title3.text_frame.paragraphs[0].runs[0].font.size = Pt(50)
    title3.text_frame.paragraphs[0].runs[0].font.bold = True

    image_placeholder = slide3.placeholders[2]
    for shape in slide3.shapes:
        if shape == image_placeholder:
            slide3.shapes._spTree.remove(shape._element)

    image_placeholder3 = slide3.placeholders[1]

    image_path3 = 'images/3.jpg'
    left3 = image_placeholder3.left
    top3 = image_placeholder3.top
    width3 = image_placeholder3.width
    height3 = image_placeholder3.height

    # Calculate the aspect ratio of the image
    image_ratio3 = width3 / height3

    # Load the image to get its original dimensions
    image3 = slide3.shapes.add_picture(image_path3, left3, top3)
    original_width3 = image3.width
    original_height3 = image3.height

    # Calculate the new dimensions to fit within the placeholder
    if image_ratio3 > 1:
        new_width3 = width3
        new_height3 = int(new_width3 / image_ratio3)
    else:
        new_height3 = height3
        new_width3 = int(new_height3 * image_ratio3)

    # Calculate the offset to center the image within the placeholder
    left_offset3 = (width3 - new_width3) // 2
    top_offset3 = (height3 - new_height3) // 2

    # Remove the temporary image shape
    slide3.shapes._spTree.remove(image3._element)

    # Add the resized image to the placeholder
    slide3.shapes.add_picture(image_path3, left1 + left_offset, top1 + top_offset, new_width, new_height)

    slide4 = presentation.slides.add_slide(presentation.slide_layouts[3])
    title4 = slide4.shapes.title
    title4.text = slide_t[1]
    content4 = slide4.placeholders[1]
    content4.text = slide_c[1]

    slide4.background.fill.solid()
    slide4.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
    #
    for shape in slide4.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # RGB values for white
                    run.font.size = Pt(18)

    title4.text_frame.paragraphs[0].runs[0].font.size = Pt(40)
    title4.text_frame.paragraphs[0].runs[0].font.bold = True

    image_placeholder4 = slide4.placeholders[1]

    image_path4 = 'images/4.jpg'
    left4 = image_placeholder4.left
    top4 = image_placeholder4.top
    width4 = image_placeholder4.width
    height4 = image_placeholder4.height

    # Calculate the aspect ratio of the image
    image_ratio4 = width4 / height4

    # Load the image to get its original dimensions
    image4 = slide4.shapes.add_picture(image_path4, left4, top4)
    original_width4 = image4.width
    original_height4 = image4.height

    # Calculate the new dimensions to fit within the placeholder
    if image_ratio4 > 1:
        new_width4 = width4
        new_height4 = int(new_width4 / image_ratio4)
    else:
        new_height4 = height4
        new_width4 = int(new_height4 * image_ratio4)

    # Calculate the offset to center the image within the placeholder
    left_offset4 = (width4 - new_width4) // 2
    top_offset4 = (height4 - new_height4) // 2

    # Remove the temporary image shape
    slide4.shapes._spTree.remove(image4._element)
    slide4.shapes.add_picture(image_path4, left1 + left_offset, top1 + top_offset, new_width, new_height)

    slide5 = presentation.slides.add_slide(presentation.slide_layouts[3])
    title5 = slide5.shapes.title
    title5.text = slide_t[2]
    content5 = slide5.placeholders[1]
    content5.text = slide_c[2]
    slide5.background.fill.solid()
    slide5.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
    #
    for shape in slide5.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # RGB values for white
                    run.font.size = Pt(18)

    title5.text_frame.paragraphs[0].runs[0].font.size = Pt(40)
    title5.text_frame.paragraphs[0].runs[0].font.bold = True

    image_placeholder5 = slide5.placeholders[1]

    image_path5 = 'images/5.jpg'
    left5 = image_placeholder5.left
    top5 = image_placeholder5.top
    width5 = image_placeholder5.width
    height5 = image_placeholder5.height

    # Calculate the aspect ratio of the image
    image_ratio5 = width5 / height5

    # Load the image to get its original dimensions
    image5 = slide5.shapes.add_picture(image_path5, left5, top5)
    original_width5 = image5.width
    original_height5 = image5.height

    # Calculate the new dimensions to fit within the placeholder
    if image_ratio5 > 1:
        new_width5 = width5
        new_height5 = int(new_width5 / image_ratio5)
    else:
        new_height5 = height5
        new_width5 = int(new_height5 * image_ratio5)

    # Calculate the offset to center the image within the placeholder
    left_offset5 = (width5 - new_width5) // 2
    top_offset5 = (height5 - new_height5) // 2

    # Remove the temporary image shape
    slide5.shapes._spTree.remove(image5._element)
    slide5.shapes.add_picture(image_path5, left1 + left_offset, top1 + top_offset, new_width, new_height)

    slide7 = presentation.slides.add_slide(presentation.slide_layouts[3])
    title7 = slide7.shapes.title
    title7.text = slide_t[3]
    content7 = slide7.placeholders[1]
    content7.text = slide_c[3]

    slide7.background.fill.solid()
    slide7.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
    #
    for shape in slide7.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # RGB values for white
                    run.font.size = Pt(18)

    title7.text_frame.paragraphs[0].runs[0].font.size = Pt(40)
    title7.text_frame.paragraphs[0].runs[0].font.bold = True

    image_placeholder7 = slide7.placeholders[1]

    image_path7 = 'images/7.jpg'
    left7 = image_placeholder7.left
    top7 = image_placeholder7.top
    width7 = image_placeholder7.width
    height7 = image_placeholder7.height

    # Calculate the aspect ratio of the image
    image_ratio7 = width7 / height7

    # Load the image to get its original dimensions
    image7 = slide7.shapes.add_picture(image_path7, left7, top7)
    original_width7 = image7.width
    original_height7 = image7.height

    # Calculate the new dimensions to fit within the placeholder
    if image_ratio7 > 1:
        new_width7 = width7
        new_height7 = int(new_width7 / image_ratio7)
    else:
        new_height7 = height7
        new_width7 = int(new_height7 * image_ratio7)

    # Calculate the offset to center the image within the placeholder
    left_offset7 = (width7 - new_width7) // 2
    top_offset7 = (height7 - new_height7) // 2

    # Remove the temporary image shape
    slide7.shapes._spTree.remove(image7._element)
    slide7.shapes.add_picture(image_path7, left1 + left_offset, top1 + top_offset, new_width, new_height)

    slide8 = presentation.slides.add_slide(presentation.slide_layouts[3])
    title8 = slide8.shapes.title
    title8.text = slide_t[4]
    content8 = slide8.placeholders[1]
    content8.text = slide_c[4]
    slide8.background.fill.solid()
    slide8.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
    #
    for shape in slide8.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # RGB values for white
                    run.font.size = Pt(18)

    title8.text_frame.paragraphs[0].runs[0].font.size = Pt(40)
    title8.text_frame.paragraphs[0].runs[0].font.bold = True

    image_placeholder8 = slide8.placeholders[1]

    image_path8 = 'images/8.jpg'
    left8 = image_placeholder8.left
    top8 = image_placeholder8.top
    width8 = image_placeholder8.width
    height8 = image_placeholder8.height

    # Calculate the aspect ratio of the image
    image_ratio8 = width8 / height8

    # Load the image to get its original dimensions
    image8 = slide8.shapes.add_picture(image_path8, left8, top8)
    original_width8 = image8.width
    original_height8 = image8.height

    # Calculate the new dimensions to fit within the placeholder
    if image_ratio8 > 1:
        new_width8 = width8
        new_height8 = int(new_width8 / image_ratio8)
    else:
        new_height8 = height8
        new_width8 = int(new_height8 * image_ratio8)

    # Calculate the offset to center the image within the placeholder
    left_offset8 = (width8 - new_width8) // 2
    top_offset8 = (height8 - new_height8) // 2

    # Remove the temporary image shape
    slide8.shapes._spTree.remove(image8._element)
    slide8.shapes.add_picture(image_path8, left1 + left_offset, top1 + top_offset, new_width, new_height)

    # Save the presentation
    presentation.save('test.pptx')

    # Reset the file pointer of the BytesIO object to the beginning
    # ppt_bytes.seek(0)

    return ppt_bytes


get_presentation('Robots')